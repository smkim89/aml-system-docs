"""
BO 와이어프레임 도형 라이브러리 (정본: docs/plan/sample.pptx)
=================================================================
핵심 원칙: 와이어프레임은 **실제 도형(rect)** 으로 그린다. ASCII 박스 문자
(┌─┐│└┘)를 모노스페이스 텍스트로 박아넣지 않는다.

팔레트(Ant Design, sample.pptx 추출):
  텍스트 262626 / 보조 595959·8C8C8C / 보더 D9D9D9 / 패널 FAFAFA
  프라이머리 1890FF / 성공 52C41A / 경고 FFA940
  콜아웃 배경 FFF7E6 · 텍스트 874D00 / 활성네비 E6F7FF / 헤더 001529·검색 003A6C
폰트: 맑은 고딕.  레이아웃: 13.33 x 7.5in.
  네비 x0.11 w1.86 / 콘텐츠 x2.13 w7.49 / 정보패널 x9.76 w3.28
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

FONT = "맑은 고딕"
C = {
    "text": "262626", "sub": "595959", "mute": "8C8C8C", "border": "D9D9D9",
    "panel": "FAFAFA", "blue": "1890FF", "green": "52C41A", "orange": "FFA940",
    "callbg": "FFF7E6", "calltx": "874D00", "active": "E6F7FF",
    "header": "001529", "search": "003A6C", "white": "FFFFFF", "redtx": "CF1322",
}
# 레이아웃 앵커
NAV_X, NAV_W = 0.11, 1.86
CON_X, CON_W = 2.13, 7.49
INFO_X, INFO_W = 9.76, 3.28
TOP_Y = 0.85          # 헤더바 상단
CON_TOP = 1.86        # 콘텐츠 본문 시작(브레드크럼/타이틀 아래)


def _rgb(h): return RGBColor.from_string(h)


def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(line_w)
    return sp


def text(slide, x, y, w, h, runs, size=10, color="262626", bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, mono_off=True, wrap=True):
    """runs: str 또는 [(s,size,color,bold), ...] 또는 [[run,...]=문단]."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [[(runs, size, color, bold)]]
    elif runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (s, sz, col, bd) in para:
            r = p.add_run(); r.text = s
            r.font.name = s_(r); r.font.size = Pt(sz); r.font.bold = bd
            r.font.color.rgb = _rgb(col)
    return tb


def s_(run):
    """폰트를 맑은 고딕으로(라틴/동아시아 모두)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)
    return FONT


# ── 프레임 컴포넌트 ─────────────────────────────────────────────
def page_title(slide, top, screen_id):
    text(slide, 0.13, 0.14, 9.4, 0.40, [(top, 18, C["text"], True)])
    text(slide, 0.13, 0.57, 9.4, 0.24, [(screen_id, 10, C["mute"], False)])


def header_bar(slide, search_ph="검색...", admin="관리자 admin ▼", brand="HANPASS", action=None):
    rect(slide, NAV_X, TOP_Y, NAV_W + CON_W + 0.0, 0.33, C["header"])
    text(slide, NAV_X + 0.11, TOP_Y + 0.08, 1.3, 0.17, [(brand, 11, C["white"], True)])
    rect(slide, CON_X, TOP_Y + 0.06, 3.83, 0.22, C["search"])
    text(slide, CON_X + 0.14, TOP_Y + 0.09, 3.5, 0.16, [(search_ph, 9, "BFBFBF", False)])
    right = NAV_X + NAV_W + CON_W
    text(slide, right - 1.45, TOP_Y + 0.09, 1.4, 0.16,
         [(admin, 9, C["white"], False)], align=PP_ALIGN.RIGHT)
    if action:
        aw = max(0.95, 0.30 + len(action) * 0.10)
        ax = right - 1.45 - 0.14 - aw
        rect(slide, ax, TOP_Y + 0.055, aw, 0.22, C["blue"])
        text(slide, ax, TOP_Y + 0.09, aw, 0.16, [(action, 8.5, C["white"], True)],
             align=PP_ALIGN.CENTER)


def nav_panel(slide, items, active=0):
    rect(slide, NAV_X, 1.18, NAV_W, 6.12, C["panel"], line=C["border"], line_w=0.5)
    y0, step = 1.29, 0.255
    if 0 <= active < len(items):
        rect(slide, NAV_X, y0 + active * step - 0.03, NAV_W, 0.23, C["active"])
    for i, it in enumerate(items):
        col = C["blue"] if i == active else C["text"]
        text(slide, NAV_X + 0.11, y0 + i * step, NAV_W - 0.16, 0.17,
             [("● " + it, 9.5, col, i == active)])


def nav_tree(slide, areas, active_key=None):
    """2영역(운영/설정)·3단(영역→기능그룹→메뉴) NAV.
    areas: [(area_label, [(group_label, [(menu_label, key), ...]), ...]), ...]
    active_key 와 일치하는 menu key 행을 하이라이트하고 active 색을 입힌다.
    행 수에 따라 step·폰트를 동적 축소해 6.12in 패널 안에 항상 수용."""
    rect(slide, NAV_X, 1.18, NAV_W, 6.12, C["panel"], line=C["border"], line_w=0.5)
    rows = 0
    for _, groups in areas:
        rows += 1
        for _, menus in groups:
            rows += 1 + len(menus)
    step = min(0.255, 6.00 / max(rows, 1))
    fs = 8.0 if rows > 24 else 9.0
    y = 1.27
    for area_label, groups in areas:
        rect(slide, NAV_X, y - 0.02, NAV_W, step - 0.01, C["blue"])
        text(slide, NAV_X + 0.08, y, NAV_W - 0.12, step - 0.02,
             [("▼ " + area_label, fs + 1.0, C["white"], True)], wrap=False)
        y += step
        for group_label, menus in groups:
            text(slide, NAV_X + 0.12, y, NAV_W - 0.18, step - 0.02,
                 [("▸ " + group_label, fs, C["mute"], True)], wrap=False)
            y += step
            for menu_label, key in menus:
                on = (key == active_key)
                if on:
                    rect(slide, NAV_X, y - 0.01, NAV_W, step - 0.01, C["active"])
                col = C["blue"] if on else C["text"]
                text(slide, NAV_X + 0.28, y, NAV_W - 0.34, step - 0.02,
                     [(menu_label, fs, col, on)], wrap=False)
                y += step


def breadcrumb_title(slide, crumb, title):
    text(slide, CON_X, 1.24, CON_W, 0.16, [(crumb, 9, C["mute"], False)])
    text(slide, CON_X, 1.42, CON_W, 0.35, [(title, 16, C["text"], True)])


def info_panel(slide, screen_id, bullets):
    rect(slide, INFO_X, TOP_Y, INFO_W, 6.45, C["white"], line=C["border"], line_w=0.75)
    text(slide, INFO_X + 0.07, TOP_Y + 0.09, INFO_W - 0.14, 0.20,
         [("기능 설명  |  " + screen_id, 10.5, C["text"], True)])
    rect(slide, INFO_X + 0.07, TOP_Y + 0.34, INFO_W - 0.16, 0.02, C["blue"])
    y = TOP_Y + 0.43
    for b in bullets:
        lead = b[0] if b[:1] in ("•", "▸") else "•"
        body = b[1:].strip() if b[:1] in ("•", "▸") else b
        # 대략 줄 수 추정 → 높이
        import math
        lines = max(1, math.ceil(len(body) / 18))
        h = 0.18 * lines + 0.04
        text(slide, INFO_X + 0.07, y, INFO_W - 0.16, h,
             [[(lead + " ", 9, C["blue"], True), (body, 9, C["sub"], False)]])
        y += h + 0.05


# ── 본문 블록 컴포넌트 (y-커서 기반) ───────────────────────────
def filters(slide, y, items):
    text(slide, CON_X, y, CON_W, 0.16, [("필터", 9, C["mute"], False)])
    x = CON_X
    for it in items:
        w = max(1.1, 0.30 + len(it) * 0.105)
        rect(slide, x, y + 0.18, w, 0.24, C["white"], line=C["border"], line_w=0.5)
        text(slide, x + 0.08, y + 0.225, w - 0.10, 0.16, [(it + "  ▼", 9, C["text"], False)],
             wrap=False)
        x += w + 0.12
    return y + 0.54


def kpi_cards(slide, y, cards):
    """cards: [(label, number, sub, accent_key), ...]  최대 4장."""
    n = len(cards); gap = 0.12
    w = (CON_W - gap * (n - 1)) / n
    accents = {"blue": C["blue"], "green": C["green"], "orange": C["orange"], "red": C["redtx"]}
    for i, (label, num, sub, acc) in enumerate(cards):
        x = CON_X + i * (w + gap)
        rect(slide, x, y, w, 0.98, C["white"], line=C["border"], line_w=0.5)
        rect(slide, x, y, w, 0.07, accents.get(acc, C["blue"]))
        text(slide, x + 0.12, y + 0.17, w - 0.24, 0.18, [(label, 9.5, C["sub"], False)])
        text(slide, x + 0.12, y + 0.40, w - 0.24, 0.32, [(num, 19, C["text"], True)])
        text(slide, x + 0.12, y + 0.74, w - 0.24, 0.18, [(sub, 8, C["mute"], False)])
    return y + 0.98 + 0.16


def callout(slide, y, title, lines, h=None):
    h = h or (0.30 + 0.20 * len(lines))
    rect(slide, CON_X, y, CON_W, h, C["callbg"], line=C["orange"], line_w=1.0)
    text(slide, CON_X + 0.14, y + 0.10, CON_W - 0.28, 0.18, [("[ " + title + " ]", 9.5, C["calltx"], True)])
    yy = y + 0.32
    for ln in lines:
        text(slide, CON_X + 0.20, yy, CON_W - 0.36, 0.18, [(ln, 8.5, C["calltx"], False)])
        yy += 0.20
    return y + h + 0.14


def entry_banner(slide, y, trigger):
    """드릴다운/상세 화면 상단 '진입 경로(트리거)' 배너 — 어떤 화면의 어떤 버튼/행을 눌러 왔는지 명시.
    화면 ID 간 흐름 단절을 막는다. trigger 예: 'AML-RA-001 ② 고위험 목록에서 [대상 행 ▶] 클릭'."""
    rect(slide, CON_X, y, CON_W, 0.30, C["active"], line=C["blue"], line_w=0.8)
    text(slide, CON_X + 0.12, y + 0.055, CON_W - 0.24, 0.19,
         [("↩ 진입 경로  ", 9, C["blue"], True), (trigger, 9, C["text"], False)], wrap=False)
    return y + 0.30 + 0.10


def panel_table(slide, x, y, w, h, title, headers, rows, col_w=None, wrap=True):
    """경계 있는 패널 + 표(헤더 + 행). col_w: 비율 리스트.
    wrap=False면 셀을 한 줄로 렌더(줄바꿈으로 다음 행을 침범하지 않음)."""
    rect(slide, x, y, w, h, C["white"], line=C["border"], line_w=0.5)
    text(slide, x + 0.12, y + 0.10, w - 0.24, 0.18, [(title, 10, C["text"], True)])
    ty = y + 0.40
    inner = w - 0.24
    if not col_w:
        col_w = [1.0 / len(headers)] * len(headers)
    cx = x + 0.12
    # 헤더
    if headers:
        xx = cx
        for j, hd in enumerate(headers):
            text(slide, xx, ty, inner * col_w[j], 0.18, [(hd, 8.5, C["mute"], True)])
            xx += inner * col_w[j]
        ty += 0.24
        rect(slide, x + 0.12, ty - 0.04, inner, 0.008, C["border"])
    # 행
    for row in rows:
        xx = cx
        for j, cell in enumerate(row):
            col = C["blue"] if (j == len(row) - 1 and str(cell).strip().lstrip("-").isdigit()) else C["text"]
            text(slide, xx, ty, inner * col_w[j], 0.18, [(str(cell), 8.5, col, False)],
                 align=(PP_ALIGN.RIGHT if (j == len(row) - 1 and str(cell).replace(",", "").lstrip("-").isdigit()) else PP_ALIGN.LEFT),
                 wrap=wrap)
            xx += inner * col_w[j]
        ty += 0.225
    return y + h + 0.14


def table_block(slide, y, h, title, headers, rows, col_w=None):
    return panel_table(slide, CON_X, y, CON_W, h, title, headers, rows, col_w)


def two_panels(slide, y, h, left, right):
    """left/right: (title, headers, rows, col_w)."""
    gap = 0.15; w = (CON_W - gap) / 2
    panel_table(slide, CON_X, y, w, h, *left)
    panel_table(slide, CON_X + w + gap, y, w, h, *right)
    return y + h + 0.14


def tab_chips(slide, y, tabs, active=0):
    """탭 칩 스트립(상세/폼 화면 상단). active 탭만 파란 강조."""
    x = CON_X
    for i, t in enumerate(tabs):
        w = max(1.0, 0.30 + len(t) * 0.13)
        on = (i == active)
        rect(slide, x, y, w, 0.26, C["active"] if on else C["white"],
             line=(C["blue"] if on else C["border"]), line_w=(0.9 if on else 0.5))
        text(slide, x + 0.08, y + 0.045, w - 0.14, 0.18,
             [(t, 9, C["blue"] if on else C["sub"], on)], wrap=False)
        x += w + 0.10
    return y + 0.26 + 0.16


def form_panel(slide, x, y, w, h, title, fields, btns=None):
    """라벨 + 입력 박스 폼(마스터 생성/수정).
    fields: [(label, value_or_placeholder, kind), ...]
      kind: 'input'(텍스트박스) / 'radio'(라디오 표시) / 'check'(체크) / 'text'(읽기전용)
    btns: ["취소","등록"] 우하단 버튼."""
    rect(slide, x, y, w, h, C["white"], line=C["border"], line_w=0.5)
    text(slide, x + 0.14, y + 0.12, w - 0.28, 0.20, [(title, 10.5, C["text"], True)])
    rect(slide, x + 0.14, y + 0.40, w - 0.28, 0.014, C["blue"])
    ry = y + 0.54
    lab_w = 1.55
    inp_x = x + 0.14 + lab_w + 0.10
    inp_w = (x + w) - inp_x - 0.16
    for (lab, val, kind) in fields:
        text(slide, x + 0.14, ry + 0.04, lab_w, 0.20, [(lab, 9, C["sub"], False)])
        if kind == "input":
            rect(slide, inp_x, ry, inp_w, 0.26, C["white"], line=C["border"], line_w=0.5)
            text(slide, inp_x + 0.08, ry + 0.045, inp_w - 0.14, 0.18,
                 [(val, 9, C["mute"], False)])
        elif kind == "radio":
            text(slide, inp_x, ry + 0.04, inp_w, 0.20, [(val, 9, C["text"], False)])
        elif kind == "check":
            rect(slide, inp_x, ry + 0.02, 0.18, 0.18, C["white"], line=C["blue"], line_w=0.9)
            text(slide, inp_x + 0.10, ry + 0.04, 0.10, 0.18, [("✓", 8, C["blue"], True)])
            text(slide, inp_x + 0.28, ry + 0.04, inp_w - 0.30, 0.20, [(val, 9, C["text"], False)])
        else:  # text
            text(slide, inp_x, ry + 0.04, inp_w, 0.20, [(val, 9, C["text"], False)])
        ry += 0.40
    if btns:
        bx = x + w - 0.16
        for b in reversed(btns):
            bw = max(0.80, 0.30 + len(b) * 0.11)
            bx -= bw
            prim = (b == btns[-1])
            rect(slide, bx, y + h - 0.42, bw, 0.28,
                 C["blue"] if prim else C["white"], line=C["border"], line_w=0.5)
            text(slide, bx, y + h - 0.375, bw, 0.18,
                 [(b, 9, C["white"] if prim else C["text"], prim)], align=PP_ALIGN.CENTER)
            bx -= 0.10
    return y + h + 0.14


def form_block(slide, y, h, title, fields, btns=None):
    return form_panel(slide, CON_X, y, CON_W, h, title, fields, btns)


def cover_slide(p, title, subtitle, meta_lines, brand="HANPASS  ·  SaaS FDS Platform"):
    """1번 커버 슬라이드. brand 기본값은 FDS 호환(기존 호출 불변)."""
    s = add_slide(p)
    rect(s, 0, 0, 13.333, 7.5, C["header"])
    rect(s, 0, 3.05, 13.333, 0.02, C["blue"])
    text(s, 1.0, 1.05, 11.3, 0.4, [(brand, 14, "8CC8FF", True)])
    text(s, 1.0, 2.0, 11.3, 1.0, [(title, 34, C["white"], True)])
    text(s, 1.0, 3.25, 11.3, 0.5, [(subtitle, 15, "BFBFBF", False)])
    yy = 4.6
    for ln in meta_lines:
        text(s, 1.0, yy, 11.3, 0.3, [[("•  ", 11, C["blue"], True), (ln, 11, "D9D9D9", False)]])
        yy += 0.42
    return s


def history_slide(p, title, headers, rows, col_w=None, clip=None):
    """2번 변경 이력 슬라이드(전폭 표).
    '변경 내역'(마지막 열)은 한 줄에 맞게 절단해 행 간 겹침을 방지한다.
    clip: 마지막 열 최대 글자 수(기본=열 너비 기반 자동). 상세 이력은 PRD 마크다운에 보존."""
    s = add_slide(p)
    text(s, 0.55, 0.45, 12.0, 0.5, [(title, 24, C["text"], True)])
    rect(s, 0.55, 1.05, 12.23, 0.02, C["blue"])
    # 마지막 열 절단(한 줄 유지) — 한글 폭 고려 보수적 한도
    if clip is None:
        last_ratio = (col_w[-1] if col_w else 1.0 / max(1, len(headers)))
        clip = max(20, int((12.23 - 0.24) * last_ratio / 0.095))
    def _clip(s_):
        s_ = str(s_)
        return s_ if len(s_) <= clip else s_[:clip - 1].rstrip() + "…"
    rows = [list(r[:-1]) + [_clip(r[-1])] for r in rows]
    h = min(5.7, 0.5 + 0.30 * (len(rows) + 1))
    panel_table(s, 0.55, 1.35, 12.23, h, "변경 이력", headers, rows, col_w, wrap=False)
    return s


def condition_builder(slide, y, h, title, combine, conditions, footer_btns=None,
                      group_note=None):
    """추가 조건 빌더(여러 조건을 AND/OR 결합).
    combine: 'AND'(모두 만족) 또는 'OR'(하나라도)
    conditions: [(필드, 연산자, 값), ...]  각 조건 = 필드+연산자+값
    각 조건 행 사이에 AND/OR 칩, 하단에 [+ 조건/그룹] 버튼·그룹 결합 표기."""
    x, w = CON_X, CON_W
    rect(slide, x, y, w, h, C["white"], line=C["border"], line_w=0.5)
    text(slide, x + 0.14, y + 0.11, w - 0.28, 0.20, [(title, 10.5, C["text"], True)])
    rect(slide, x + 0.14, y + 0.38, w - 0.28, 0.014, C["blue"])
    # 결합 방식 라디오
    cy = y + 0.48
    text(slide, x + 0.14, cy + 0.03, 1.0, 0.18, [("결합 방식", 9, C["sub"], False)])
    a_on, o_on = (combine == "AND"), (combine == "OR")
    text(slide, x + 1.10, cy + 0.03, 2.2, 0.18,
         [("●" if a_on else "○", 9, C["blue"] if a_on else C["mute"], True),
          (" 모두 만족(AND)", 9, C["text"], a_on)])
    text(slide, x + 3.20, cy + 0.03, 2.2, 0.18,
         [("●" if o_on else "○", 9, C["blue"] if o_on else C["mute"], True),
          (" 하나라도(OR)", 9, C["text"], o_on)])
    # 조건 행
    inner_x, inner_w = x + 0.14, w - 0.28
    fw, ow, vw = inner_w * 0.40, inner_w * 0.22, inner_w * 0.26
    ry = cy + 0.26
    for i, (fld, op, val) in enumerate(conditions):
        if i > 0:
            # AND/OR 결합 칩
            rect(slide, inner_x, ry, 0.62, 0.16, C["active"], line=C["blue"], line_w=0.6)
            text(slide, inner_x, ry + 0.005, 0.62, 0.15, [(combine, 8, C["blue"], True)],
                 align=PP_ALIGN.CENTER)
            ry += 0.205
        rect(slide, inner_x, ry, inner_w, 0.30, C["panel"], line=C["border"], line_w=0.4)
        text(slide, inner_x + 0.06, ry + 0.075, 0.50, 0.16, [("조건 " + str(i + 1), 8, C["mute"], True)])
        bx = inner_x + 0.62
        for (cap, bw) in ((fld + "  ▼", fw), (op + "  ▼", ow), (val, vw)):
            rect(slide, bx, ry + 0.045, bw - 0.08, 0.21, C["white"], line=C["border"], line_w=0.4)
            text(slide, bx + 0.06, ry + 0.075, bw - 0.18, 0.16, [(cap, 8.5, C["text"], False)], wrap=False)
            bx += bw
        text(slide, inner_x + inner_w - 0.18, ry + 0.075, 0.16, 0.16, [("✕", 8, C["mute"], False)])
        ry += 0.30
    # 하단 버튼 + 그룹 결합
    fy = y + h - 0.34
    bx = inner_x
    for b in (footer_btns or ["+ 조건 추가", "+ 그룹(괄호) 추가"]):
        bw = max(0.95, 0.30 + len(b) * 0.085)
        rect(slide, bx, fy, bw, 0.24, C["white"], line=C["blue"], line_w=0.6)
        text(slide, bx, fy + 0.04, bw, 0.16, [(b, 8.5, C["blue"], True)], align=PP_ALIGN.CENTER)
        bx += bw + 0.10
    if group_note:
        text(slide, inner_x + inner_w - 2.2, fy + 0.04, 2.2, 0.16,
             [(group_note, 8.5, C["sub"], False)], align=PP_ALIGN.RIGHT)
    return y + h + 0.14


def new_deck():
    p = Presentation()
    p.slide_width = Inches(13.333); p.slide_height = Inches(7.5)
    return p


def add_slide(p):
    return p.slides.add_slide(p.slide_layouts[6])  # blank
